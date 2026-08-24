#!/usr/bin/env python3
"""Ce que l'ingestion RAG sait lire — et ce qu'elle refuse proprement.

**L'ÉTAT MESURÉ AVANT D'AJOUTER QUOI QUE CE SOIT**, parce que « c'est
probablement déjà couvert » n'est pas une réponse :

======================  =========================================================
format                  comportement du 2026-08-24, avant ce lot
======================  =========================================================
``.pdf``                lu (pypdf), déjà
``.docx``               **lu** — python-docx ne servait donc PAS qu'à l'écriture.
                        Mais ``doc.paragraphs`` n'inclut pas les cellules de
                        tableau : un document d'un paragraphe et d'un tableau à
                        deux cellules rendait **14 caractères**, le tableau perdu
                        en silence.
``.pptx`` / ``.xlsx``   chaîne VIDE, sans un mot. ``index_file`` sortait tôt, donc
                        un fichier accepté à l'upload s'indexait à zéro chunk.
inconnu (``.odt``)      chaîne vide aussi — indiscernable des deux précédents
======================  =========================================================

Ce que ces tests gardent, dans cet ordre d'importance :

1. **le texte extrait correspond à la source** — pas « non vide », ce qui
   passerait sur un extracteur qui rendrait le nom du fichier ;
2. **le contenu des tableaux d'un .docx est là** — c'est la régression qu'on
   corrige, et un test qui n'assemble qu'un paragraphe ne la verrait pas ;
3. **un format non supporté est refusé, pas planté** — et la frontière est à
   l'upload (400 avec la liste), pas dans l'extracteur ;
4. **une seule liste d'extensions** : elle était écrite trois fois côté backend.

Les fichiers sont **fabriqués ici**, pas versionnés : un .pptx de référence dans
le dépôt serait un binaire opaque que personne ne relit, et dont personne ne
saurait dire ce qu'il est censé contenir.

Usage :
    python test_ingestion_documents.py
"""

import importlib.util
import os
import shutil
import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import _test_env  # noqa: F401  — isole les chemins AVANT tout import de core.*

from core.rag import SUPPORTED_EXTENSIONS, RAGEngine  # noqa: E402


def _present(module: str) -> bool:
    """Le lecteur est-il installé ?

    Les trois sont déclarés dans `requirements.txt` ET dans le jeu minimal de
    `ci.yml` — donc ces gardes ne devraient jamais s'activer. Elles sont là pour
    qu'un environnement partiel donne un test SAUTÉ et nommé plutôt qu'une erreur
    d'import qui ferait tomber tout le fichier, y compris les tests qui ne
    dépendent d'aucun lecteur (la liste d'extensions, le refus des formats
    inconnus).
    """
    try:
        return importlib.util.find_spec(module) is not None
    except (ImportError, ValueError):
        return False


class _Fichiers(unittest.TestCase):
    """Base : un dossier temporaire par classe, nettoyé à la fin."""

    @classmethod
    def setUpClass(cls):
        cls.dossier = Path(tempfile.mkdtemp(prefix="epure-test-ingestion-"))

    @classmethod
    def tearDownClass(cls):
        shutil.rmtree(cls.dossier, ignore_errors=True)


@unittest.skipUnless(_present("docx"), "python-docx absent")
class DocxTest(_Fichiers):
    """.docx : paragraphes ET tableaux.

    Le tableau est le sujet. Avant ce lot, `_extract_text_from_path` rendait les
    paragraphes seuls — et un tableau est souvent ce qu'un document de cours
    contient de plus dense (dates, correspondances, formules).
    """

    @classmethod
    def setUpClass(cls):
        super().setUpClass()
        from docx import Document
        doc = Document()
        doc.add_paragraph("Le théorème de Thalès s'énonce ainsi.")
        doc.add_paragraph("")  # paragraphe vide : ne doit pas polluer la sortie
        doc.add_paragraph("Second paragraphe avec un accent : é.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Grandeur"
        table.cell(0, 1).text = "Unité"
        table.cell(1, 0).text = "Longueur"
        table.cell(1, 1).text = "mètre"
        cls.chemin = cls.dossier / "cours.docx"
        doc.save(cls.chemin)
        cls.texte = RAGEngine.read_file_text(str(cls.chemin))

    def test_les_paragraphes_sont_la(self):
        self.assertIn("théorème de Thalès", self.texte)
        self.assertIn("Second paragraphe avec un accent : é.", self.texte)

    def test_le_contenu_du_tableau_est_la(self):
        """LA régression corrigée. `doc.paragraphs` ne voit pas les cellules."""
        for attendu in ("Grandeur", "Unité", "Longueur", "mètre"):
            self.assertIn(attendu, self.texte, attendu)

    def test_une_ligne_de_tableau_reste_une_ligne(self):
        """Les cellules d'une même ligne sont sur la même ligne de texte.

        Sinon « Longueur » et « mètre » peuvent tomber dans deux chunks
        différents, et la correspondance que le tableau exprime est perdue alors
        que les deux mots sont indexés.
        """
        self.assertIn("Longueur\tmètre", self.texte)
        self.assertIn("Grandeur\tUnité", self.texte)

    def test_les_paragraphes_vides_ne_polluent_pas(self):
        self.assertNotIn("\n\n\n", self.texte)


@unittest.skipUnless(_present("pptx"), "python-pptx absent")
class PptxTest(_Fichiers):
    """.pptx : titres, corps, tableaux et notes du présentateur."""

    @classmethod
    def setUpClass(cls):
        super().setUpClass()
        from pptx import Presentation
        from pptx.util import Inches
        pres = Presentation()
        # Disposition 1 = titre + contenu, celle de toute présentation ordinaire.
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        slide.shapes.title.text = "Les lois de Newton"
        slide.placeholders[1].text = "Principe d'inertie\nRelation fondamentale"
        slide.notes_slide.notes_text_frame.text = (
            "Insister sur le référentiel galiléen."
        )
        # Deuxième diapositive : un tableau, qu'aucun text_frame ne couvre.
        vide = pres.slides.add_slide(pres.slide_layouts[6])
        forme = vide.shapes.add_table(2, 2, Inches(1), Inches(1),
                                      Inches(4), Inches(1))
        table = forme.table
        table.cell(0, 0).text = "Force"
        table.cell(0, 1).text = "Newton"
        table.cell(1, 0).text = "Masse"
        table.cell(1, 1).text = "kilogramme"
        cls.chemin = cls.dossier / "presentation.pptx"
        pres.save(cls.chemin)
        cls.texte = RAGEngine.read_file_text(str(cls.chemin))

    def test_le_titre_et_le_corps_sont_la(self):
        self.assertIn("Les lois de Newton", self.texte)
        self.assertIn("Principe d'inertie", self.texte)
        self.assertIn("Relation fondamentale", self.texte)

    def test_les_notes_du_presentateur_sont_la(self):
        """Souvent la seule prose complète d'une présentation en puces."""
        self.assertIn("référentiel galiléen", self.texte)
        # Préfixées, pour qu'un extrait ne se lise pas comme du contenu projeté.
        self.assertIn("Notes : Insister sur le référentiel galiléen.", self.texte)

    def test_le_tableau_d_une_diapositive_est_la(self):
        for attendu in ("Force", "Newton", "Masse", "kilogramme"):
            self.assertIn(attendu, self.texte, attendu)
        self.assertIn("Masse\tkilogramme", self.texte)

    def test_les_diapositives_sont_numerotees(self):
        """Un extrait doit pouvoir se situer : « diapositive 2 », pas un flottant."""
        self.assertIn("Diapositive 1", self.texte)
        self.assertIn("Diapositive 2", self.texte)
        self.assertLess(self.texte.index("Diapositive 1"),
                        self.texte.index("Diapositive 2"))

    def test_une_presentation_vide_rend_une_chaine_vide(self):
        """Aucune diapositive : pas de « Diapositive 1 » fantôme."""
        from pptx import Presentation
        chemin = self.dossier / "vide.pptx"
        Presentation().save(chemin)
        self.assertEqual(RAGEngine.read_file_text(str(chemin)).strip(), "")


@unittest.skipUnless(_present("openpyxl"), "openpyxl absent")
class XlsxTest(_Fichiers):
    """.xlsx : contenu cellule par cellule, feuille par feuille."""

    @classmethod
    def setUpClass(cls):
        super().setUpClass()
        import openpyxl
        classeur = openpyxl.Workbook()
        f1 = classeur.active
        f1.title = "Notes"
        f1.append(["Matière", "Note"])
        f1.append(["Mathématiques", 18])
        f1.append([None, None])           # ligne vide : ignorée
        f1.append(["Physique", 15])
        f2 = classeur.create_sheet("Coefficients")
        f2.append(["Épreuve", "Coefficient"])
        f2.append(["Écrit", 3])
        cls.chemin = cls.dossier / "bulletin.xlsx"
        classeur.save(cls.chemin)
        classeur.close()
        cls.texte = RAGEngine.read_file_text(str(cls.chemin))

    def test_le_contenu_des_cellules_est_la(self):
        for attendu in ("Matière", "Mathématiques", "18", "Physique", "15"):
            self.assertIn(attendu, self.texte, attendu)

    def test_une_ligne_reste_une_ligne(self):
        self.assertIn("Mathématiques\t18", self.texte)

    def test_toutes_les_feuilles_sont_lues(self):
        """Une seule feuille lue serait le bug le plus facile à ne pas voir : le
        fichier de test le plus naturel n'en a qu'une.
        """
        self.assertIn("Feuille : Notes", self.texte)
        self.assertIn("Feuille : Coefficients", self.texte)
        self.assertIn("Écrit\t3", self.texte)

    def test_les_lignes_vides_sont_ignorees(self):
        self.assertNotIn("\t\t", self.texte)
        self.assertNotIn("\n\n\n", self.texte)

    def test_les_nombres_ne_sont_pas_perdus(self):
        """`str(None)` donnerait « None » partout, `str(18)` doit donner « 18 »."""
        self.assertNotIn("None", self.texte)

    def test_une_formule_non_calculee_ressort_vide(self):
        """La contrepartie de `data_only=True`, fixée pour qu'elle ne surprenne pas.

        `data_only=True` rend la valeur MISE EN CACHE par le tableur, pas la
        source de la formule. Un classeur écrit par un script et jamais ouvert
        dans Excel n'a aucun cache : sa cellule de formule ressort donc vide.

        C'est le compromis assumé — `data_only=False` rendrait `=SOMME(A1:A2)`,
        du bruit qui n'est pas ce sur quoi une question porte. Le cas se produit
        exactement pour un fichier fabriqué (comme ici) et presque jamais pour un
        fichier que quelqu'un a rempli puis enregistré.
        """
        import openpyxl
        classeur = openpyxl.Workbook()
        feuille = classeur.active
        feuille.append(["a", 1])
        feuille.append(["b", 2])
        feuille["B3"] = "=SUM(B1:B2)"
        chemin = self.dossier / "formule.xlsx"
        classeur.save(chemin)
        classeur.close()
        texte = RAGEngine.read_file_text(str(chemin))
        self.assertIn("a\t1", texte)
        # Ni la formule, ni un résultat inventé.
        self.assertNotIn("SUM", texte)
        self.assertNotIn("=", texte)

    def test_le_fichier_n_est_pas_laisse_verrouille(self):
        """`read_only=True` garde le zip ouvert : sans `close()`, Windows refuse
        de supprimer le fichier ensuite — et une ré-indexation après modification
        échouerait sur un fichier occupé.
        """
        import openpyxl
        classeur = openpyxl.Workbook()
        classeur.active.append(["x", 1])
        chemin = self.dossier / "a-supprimer.xlsx"
        classeur.save(chemin)
        classeur.close()
        RAGEngine.read_file_text(str(chemin))
        chemin.unlink()          # lèverait PermissionError si le zip restait ouvert
        self.assertFalse(chemin.exists())


class FormatsRefusesTest(_Fichiers):
    """Un format non supporté : refusé clairement, jamais planté."""

    def test_une_extension_inconnue_rend_une_chaine_vide_sans_lever(self):
        """L'extracteur ne lève PAS sur un format inconnu — il n'est pas la
        frontière. La frontière est l'upload (test suivant) et `index_file`, qui
        sort avant d'appeler l'extracteur.
        """
        chemin = self.dossier / "document.odt"
        chemin.write_bytes(b"contenu quelconque")
        self.assertEqual(RAGEngine.read_file_text(str(chemin)), "")

    def test_les_formats_binaires_pre_2007_ne_sont_pas_acceptes(self):
        """`.doc`, `.ppt`, `.xls` : aucune des trois bibliothèques ne les lit.

        Les ajouter par ressemblance donnerait une erreur à l'ouverture au lieu
        d'un refus honnête à l'upload — c'est exactement le genre d'ajout « par
        famille » qu'il faut refuser.
        """
        for ext in (".doc", ".ppt", ".xls"):
            self.assertNotIn(ext, SUPPORTED_EXTENSIONS, ext)

    def test_un_fichier_corrompu_leve_au_lieu_de_mentir(self):
        """Un .docx qui n'en est pas un doit ÉCHOUER, pas rendre du vide.

        La distinction compte : une chaîne vide se confond avec « document sans
        texte » et s'indexe en silence à zéro chunk. Une exception remonte à
        l'appelant, qui sait dire « impossible de lire ce fichier ». C'est ce que
        `.pdf` fait depuis toujours ; les nouveaux formats s'alignent.
        """
        for nom, contenu in (("faux.docx", b"pas un docx"),
                             ("faux.pptx", b"pas un pptx"),
                             ("faux.xlsx", b"pas un xlsx")):
            if not _present({"faux.docx": "docx", "faux.pptx": "pptx",
                             "faux.xlsx": "openpyxl"}[nom]):
                continue
            chemin = self.dossier / nom
            chemin.write_bytes(contenu)
            with self.assertRaises(Exception, msg=nom):
                RAGEngine.read_file_text(str(chemin))

    def test_index_file_ignore_une_extension_non_supportee(self):
        """`index_file` sort AVANT l'extracteur et avant tout accès au store.

        C'est ce qui rend ce test possible sans construire de moteur : aucun
        embedding n'est chargé pour refuser un fichier.
        """
        chemin = self.dossier / "ignore.odt"
        chemin.write_bytes(b"x")
        moteur = object.__new__(RAGEngine)   # pas d'__init__ : ni torch ni store
        self.assertIsNone(RAGEngine.index_file(moteur, str(chemin)))


class ListeUniqueTest(unittest.TestCase):
    """Une seule liste d'extensions côté backend.

    Il y en avait trois pour une notion : `core/rag.py:SUPPORTED_EXTENSIONS`, la
    copie `_SUPPORTED_EXT` du routeur Réglages, et la phrase du message d'erreur
    de l'upload qui énumérait les types à la main. Ajouter `.pptx` demandait donc
    de penser à trois endroits, et l'oubli le plus probable — la copie du routeur
    — produit le pire symptôme : un fichier refusé à l'upload alors que le moteur
    sait le lire, ou accepté alors qu'il ne sait pas, donc indexé à vide.
    """

    def test_le_routeur_partage_la_liste_du_moteur(self):
        from modules.settings.router import _SUPPORTED_EXT
        self.assertIs(_SUPPORTED_EXT, SUPPORTED_EXTENSIONS)

    def test_les_nouveaux_formats_sont_dans_la_liste(self):
        for ext in (".pptx", ".xlsx", ".docx"):
            self.assertIn(ext, SUPPORTED_EXTENSIONS, ext)

    def test_chaque_extension_de_la_liste_est_reellement_traitee(self):
        """La liste ne doit rien promettre que l'extracteur ne sait pas lire.

        Vérifié par le CODE de `_extract_text_from_path` plutôt que par un
        fichier de chaque type : les images, par exemple, y sont traitées par un
        message et non par une extraction. Ce qu'on garde ici, c'est qu'aucune
        extension n'y tombe dans le `return ""` final par oubli.
        """
        import inspect
        source = inspect.getsource(RAGEngine._extract_text_from_path)
        for ext in SUPPORTED_EXTENSIONS:
            self.assertIn(f"'{ext}'", source, ext)


if __name__ == "__main__":
    unittest.main(verbosity=2)

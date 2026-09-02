import functools
import json
import logging
import os
import threading
from datetime import datetime
from pathlib import Path
from typing import Optional

import pypdf
import yaml
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

from core.models import premier_modele_vision_disponible
from core.paths import resolve_vector_dir
from core.vector_store import VectorStore

logger = logging.getLogger(__name__)

#: Extensions que `_extract_text_from_path` sait lire. **Source unique** : le
#: routeur Réglages l'importe (`_SUPPORTED_EXT`) au lieu d'en tenir une copie, et
#: le message d'erreur de l'upload en est dérivé. Il y avait trois listes pour une
#: notion, dont deux à mettre à jour de mémoire.
#:
#: `.pptx`, `.xlsx` ajoutés le 2026-08-24. `.docx` était déjà là — et lu, contre
#: l'idée que `python-docx` ne servait qu'à l'écriture ; ce qui manquait était le
#: contenu de ses TABLEAUX (cf. `_texte_docx`).
#:
#: Ce qui n'y est PAS et ne doit pas y arriver par ressemblance : `.doc`, `.ppt`,
#: `.xls` (les formats binaires pré-2007). Aucune des trois bibliothèques ne les
#: lit — python-docx, python-pptx et openpyxl ne parlent que l'OOXML zippé — et
#: les accepter donnerait une erreur à l'ouverture au lieu d'un refus à l'upload.
SUPPORTED_EXTENSIONS = {
    '.pdf', '.docx', '.pptx', '.xlsx', '.txt', '.md', '.csv', '.json',
    '.png', '.jpg', '.jpeg', '.webp',
}

#: Sous-ensemble image de `SUPPORTED_EXTENSIONS`, utilisé par `index_file` pour
#: basculer vers `_texte_image` (modèle vision) au lieu du placeholder statique.
#: `_extract_text_from_path` garde son propre tuple littéral pour ces mêmes
#: extensions : `test_ingestion_documents.py` vérifie par lecture de son SOURCE
#: que chaque extension y apparaît en toutes lettres, une garantie qu'une
#: référence à cette constante ne satisferait pas.
_IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.webp'}

#: Nombre de lignes lues par feuille d'un classeur. Même borne que le `nrows=500`
#: du `.csv` juste en dessous, et pour la même raison : au-delà, on indexe un
#: export de base de données, pas un document qu'on lit.
_XLSX_MAX_LIGNES = 500


class RAGEngine:
    """Indexation et recherche dans les fiches surveillées.

    Le stockage est ``core.vector_store.VectorStore`` (SQLite + numpy) depuis le
    remplacement de chromadb — cf. ``docs/remplacement-vectoriel.md``. Deux
    conséquences sur cette classe, au-delà du changement de bibliothèque :

    - **Le store est INJECTÉ, plus construit ici.** ``DocAnalysisEngine`` et
      ``HistoryEngine`` travaillent sur deux autres collections du même store et
      allaient jusqu'ici le chercher dans ``rag._client``/``rag._ef`` — deux
      attributs privés d'un moteur qui n'avait aucune vocation à être un point
      de partage. Le partage est désormais explicite : ``core/runtime.py``
      construit un store et le donne aux trois. Le défaut (``store=None`` → en
      construire un) sert aux scripts et aux tests qui n'instancient que le RAG ;
      en production il n'est jamais pris.
    - **Le chemin vient de ``resolve_vector_dir()``**, plus de
      ``dirname(config.yaml)/chroma_db``. L'ancien calcul rendait l'index
      impossible à détourner : un test qui aurait construit ce moteur aurait
      écrit dans l'index réel de l'utilisateur (cf. ``core/paths.py``).
    """

    def __init__(self, config_path: str = "config.yaml", store: VectorStore | None = None,
                 llm=None):
        with open(config_path) as f:
            cfg = yaml.safe_load(f)
        rag_cfg = cfg.get("rag", {})
        self._chunk_size = rag_cfg.get("chunk_size", 500)
        self._chunk_overlap = rag_cfg.get("chunk_overlap", 50)
        self._n_results = rag_cfg.get("n_results", 3)

        self._store = store if store is not None else VectorStore(resolve_vector_dir())
        self._col = self._store.collection("fiches")

        # Injecté comme dans DocAnalysisEngine/HistoryEngine (core/runtime.py) —
        # jamais importé directement (§3.2 CLAUDE.md). Optionnel : les scripts et
        # tests légers qui construisent un RAGEngine sans vision (la majorité)
        # n'ont rien à passer, et une image s'indexe alors avec le placeholder,
        # exactement comme avant ce moteur.
        self._llm = llm

        # Per-instance LRU caches — cleared on index_file to avoid stale results
        self._query_lru = functools.lru_cache(maxsize=50)(self._do_query)
        self._query_filtered_lru = functools.lru_cache(maxsize=50)(self._do_query_filtered)

    # ── Extracteurs par format ───────────────────────────────────────────────
    #
    # Méthodes nommées et non des branches en ligne : chacune est testable seule
    # (`test_ingestion_documents.py`), sans construire de moteur ni de store —
    # ce qui compte ici, parce que construire un `RAGEngine` charge torch.
    #
    # Convention commune, calquée sur la branche `.docx` d'origine : un paquet
    # ABSENT donne un avertissement et une chaîne vide (dégradation : le paquet
    # distribué peut l'avoir perdu), un fichier ILLISIBLE laisse remonter
    # l'exception — comme le fait `.pdf` depuis toujours. Les deux ne se
    # confondent pas : l'un est une installation incomplète, l'autre un mauvais
    # fichier, et l'appelant n'a pas le même recours.

    @staticmethod
    def _texte_docx(path: str) -> str:
        """Paragraphes PUIS tableaux d'un .docx.

        `doc.paragraphs` **n'inclut pas les cellules de tableau** — mesuré : un
        document d'un paragraphe et d'un tableau à deux cellules rendait 14
        caractères, le tableau perdu en silence. Un tableau est souvent ce qu'un
        document de cours contient de plus dense (dates, formules, correspondances),
        donc c'était le contenu le plus utile qui manquait.

        Les tableaux arrivent APRÈS le corps, pas à leur place dans le fil : les
        restituer dans l'ordre demande de parcourir `doc.element.body` à la main.
        Pour de la recherche par similarité, l'ordre entre un tableau et le
        paragraphe qui l'introduit ne change rien au fait que le contenu soit
        trouvable ; l'absence, si.
        """
        try:
            from docx import Document  # python-docx
        except ImportError:
            logger.warning("python-docx non installé — pip install python-docx")
            return ""
        doc = Document(str(path))
        morceaux = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for ligne in table.rows:
                cellules = [c.text.strip() for c in ligne.cells]
                if any(cellules):
                    morceaux.append("\t".join(cellules))
        return "\n".join(morceaux)

    @staticmethod
    def _texte_pptx(path: str) -> str:
        """Texte d'un .pptx, diapositive par diapositive.

        Trois sources par diapositive, et les trois comptent :

        * les **formes à texte** — titres et corps sont l'un et l'autre des
          `shape` avec un `text_frame`, python-pptx ne les distingue pas par type.
          Les prendre tous évite de deviner un rôle de placeholder ;
        * les **tableaux** (`shape.has_table`), qu'aucun `text_frame` ne couvre ;
        * les **notes du présentateur**, souvent la seule prose complète d'une
          présentation dont les diapositives ne sont que des puces.

        Le numéro de diapositive est écrit dans le texte : c'est ce qui permet à
        une réponse de citer « diapositive 4 » au lieu d'un extrait flottant.
        """
        try:
            from pptx import Presentation  # python-pptx
        except ImportError:
            logger.warning("python-pptx non installé — pip install python-pptx")
            return ""
        presentation = Presentation(str(path))
        morceaux: list[str] = []
        for numero, slide in enumerate(presentation.slides, start=1):
            lignes: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    lignes.append(shape.text_frame.text.strip())
                if getattr(shape, "has_table", False) and shape.has_table:
                    for ligne in shape.table.rows:
                        cellules = [c.text.strip() for c in ligne.cells]
                        if any(cellules):
                            lignes.append("\t".join(cellules))
            # `has_notes_slide` avant d'y toucher : le lire crée la diapositive de
            # notes si elle n'existe pas, donc modifie l'objet en mémoire pour
            # répondre à une question de lecture.
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    lignes.append(f"Notes : {notes}")
            if lignes:
                morceaux.append(f"Diapositive {numero}\n" + "\n".join(lignes))
        return "\n\n".join(morceaux)

    @staticmethod
    def _texte_xlsx(path: str) -> str:
        """Contenu d'un .xlsx, feuille par feuille, cellule par cellule.

        Deux options d'`load_workbook` qui ne sont pas des détails :

        * ``read_only=True`` — un classeur d'export peut faire des dizaines de
          milliers de lignes, et le mode normal les charge toutes en objets ;
        * ``data_only=True`` — rend la **valeur** mise en cache des formules et
          non leur source. ``12`` plutôt que ``=SOMME(A1:A3)``, qui est ce sur
          quoi une question porte.

        **La contrepartie de ``data_only``, à connaître** : cette valeur en cache
        est écrite par le tableur, pas par le fichier. Un classeur généré par un
        script (openpyxl, pandas) et jamais ouvert dans Excel n'en a aucune —
        ses cellules de formule ressortent donc VIDES. C'est le cas le plus
        probable pour un fichier fabriqué, et le moins probable pour un fichier
        que quelqu'un a réellement rempli. Le compromis est fait dans ce sens-là ;
        `test_ingestion_documents.py` fixe le comportement pour qu'il ne surprenne
        pas plus tard.

        Le nom de la feuille est écrit dans le texte, pour la même raison que le
        numéro de diapositive : un extrait doit pouvoir se situer.
        """
        try:
            import openpyxl
        except ImportError:
            logger.warning("openpyxl non installé — pip install openpyxl")
            return ""
        classeur = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        try:
            morceaux: list[str] = []
            for feuille in classeur.worksheets:
                lignes: list[str] = []
                for n, valeurs in enumerate(feuille.iter_rows(values_only=True)):
                    if n >= _XLSX_MAX_LIGNES:
                        lignes.append(f"[… lignes suivantes non indexées "
                                      f"(limite {_XLSX_MAX_LIGNES})]")
                        break
                    cellules = ["" if v is None else str(v) for v in valeurs]
                    if any(c.strip() for c in cellules):
                        lignes.append("\t".join(cellules).rstrip("\t"))
                if lignes:
                    morceaux.append(f"Feuille : {feuille.title}\n" + "\n".join(lignes))
            return "\n\n".join(morceaux)
        finally:
            # `read_only=True` garde le zip ouvert : sans ce close, le fichier
            # reste verrouillé sous Windows, et le `rmtree` d'un test ou une
            # ré-indexation après modification échoue sur un fichier occupé.
            classeur.close()

    def _texte_image(self, path: str) -> str:
        """Description + transcription d'une image par un modèle vision.

        Remplace le placeholder muet de `_extract_text_from_path` pour les
        images : sans lui, `.png/.jpg/.jpeg/.webp` s'indexaient avec un texte
        FIXE, qui ne rendait cherchable ni un schéma ni un texte photographié —
        vérifié, aucun chemin du code n'envoyait jamais l'image à un modèle.

        Dégrade PROPREMENT vers le placeholder, jamais vers une exception — même
        convention que les extracteurs `.docx`/`.pptx`/`.xlsx` (§3.3 bis de
        CLAUDE.md : paquet/moteur absent → dégradation, fichier illisible →
        exception), et pour la même raison : un module généré ou un test léger
        qui construit ce moteur sans `llm` ne doit pas planter sur une image, et
        un modèle vision indisponible ou en panne ne doit jamais faire échouer
        l'indexation d'un fichier par ailleurs valide.
        """
        placeholder = RAGEngine._extract_text_from_path(path)
        if self._llm is None:
            logger.warning("Aucun LLM injecté dans ce RAGEngine — impossible de décrire %s, placeholder conservé", path)
            return placeholder
        modele = premier_modele_vision_disponible()
        if modele is None:
            logger.warning("Aucun modèle vision disponible pour décrire %s — placeholder conservé", path)
            return placeholder
        try:
            description = self._llm.describe_image(path, modele)
        except Exception:
            logger.exception("Échec description vision de %s (modèle %s)", path, modele)
            return placeholder
        if not description or not description.strip():
            logger.warning("Description vision vide pour %s (modèle %s) — placeholder conservé", path, modele)
            return placeholder
        return f"Image : {Path(path).name}\n\n{description.strip()}"

    @staticmethod
    def _extract_text_from_path(path: str) -> str:
        ext = Path(path).suffix.lower()
        if ext == '.pdf':
            reader = pypdf.PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == '.docx':
            return RAGEngine._texte_docx(path)
        elif ext == '.pptx':
            return RAGEngine._texte_pptx(path)
        elif ext == '.xlsx':
            return RAGEngine._texte_xlsx(path)
        elif ext in ('.txt', '.md'):
            return Path(path).read_text(encoding='utf-8', errors='ignore')
        elif ext == '.csv':
            try:
                import pandas as pd
                df = pd.read_csv(path, nrows=500)
                return df.to_string(index=False)
            except Exception:
                return Path(path).read_text(encoding='utf-8', errors='ignore')
        elif ext == '.json':
            try:
                # utf-8-sig : les JSON fournis par l'utilisateur portent souvent un BOM
                data = json.loads(Path(path).read_text(encoding='utf-8-sig'))
                return json.dumps(data, ensure_ascii=False, indent=2)[:50000]
            except Exception:
                return Path(path).read_text(encoding='utf-8', errors='ignore')
        elif ext in ('.png', '.jpg', '.jpeg', '.webp'):
            name = Path(path).name
            return f"Image : {name} (analyse vision non disponible sans modèle vision)"
        return ""

    def index_file(self, path: str) -> Optional[str]:
        """Indexe le fichier et rend le texte RÉELLEMENT indexé (ou ``None``).

        Le retour existe pour que l'appelant qui a aussi besoin du contenu —
        `_stream_load_sse` (`modules/settings/router.py`), qui construit le
        résumé affiché à l'import — réutilise CE texte au lieu de le
        recalculer avec `read_file_text`. Avant ce retour, cet appelant faisait
        une extraction séparée, et c'était plus qu'un doublon de travail pour
        les images : `read_file_text` est **statique**, donc n'appelle jamais
        `_texte_image`/le modèle vision — le résumé affiché à l'import disait
        systématiquement « je n'ai pas accès à l'image » alors que l'index
        avait la vraie description, produite juste au-dessus par le même appel
        à `index_file`. Pour les autres formats (pdf/docx/…), le contenu était
        identique des deux côtés : seule la seconde lecture/parsing était
        gaspillée, pas le résultat.

        Rend ``None`` quand rien n'a été indexé (extension non supportée,
        texte extrait vide après strip) — l'appelant ne doit pas construire de
        contenu à partir de rien.
        """
        ext = Path(path).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            return None

        full_text = (self._texte_image(path) if ext in _IMAGE_EXTENSIONS
                    else self._extract_text_from_path(path))
        if not full_text.strip():
            return None

        # Remove stale chunks before re-indexing to avoid duplicates
        try:
            self._col.delete(where={"source": str(path)})
        except Exception:
            logger.exception("Erreur suppression chunks existants pour %s", path)

        chunk_chars = self._chunk_size * 4
        overlap_chars = self._chunk_overlap * 4
        step = max(1, chunk_chars - overlap_chars)

        chunks = []
        start = 0
        while start < len(full_text):
            chunks.append(full_text[start: start + chunk_chars])
            start += step

        # mtime stocké pour la re-indexation incrémentale au démarrage (cf. watch) :
        # on saute le ré-embedding des fichiers inchangés.
        try:
            mtime = os.path.getmtime(path)
        except OSError:
            mtime = 0.0

        # `indexé_le` : quand CE passage d'indexation a eu lieu. Distinct du
        # `mtime`, qui dit quand le FICHIER a changé — deux dates différentes,
        # et déduire l'une de l'autre serait faux. Absent des chunks indexés
        # avant son ajout : l'interface affiche alors « non disponible » plutôt
        # que d'inventer.
        indexe_le = datetime.now().isoformat(timespec="seconds")
        base_id = str(path).replace("\\", "/")
        ids = [f"{base_id}::{i}" for i in range(len(chunks))]
        self._col.upsert(
            documents=chunks,
            ids=ids,
            metadatas=[{"source": str(path), "chunk": i, "mtime": mtime,
                        "indexé_le": indexe_le} for i in range(len(chunks))],
        )

        # Invalidate query caches since the index has changed
        self._query_lru.cache_clear()
        self._query_filtered_lru.cache_clear()

        return full_text

    def index_pdf(self, path: str) -> Optional[str]:
        """Backward-compat alias for index_file."""
        return self.index_file(path)

    def _do_query(self, text: str, n: int) -> str:
        count = self._col.count()
        if count == 0:
            return ""
        results = self._col.query(query_texts=[text], n_results=min(n, count))
        docs = results.get("documents", [[]])[0]
        return "\n\n---\n\n".join(d for d in docs if d)

    def _do_query_filtered(self, text: str, paths_key: tuple, n: int) -> str:
        paths = list(paths_key)
        if not paths:
            return ""
        try:
            existing = self._col.get(
                where={"source": {"$in": paths}}, include=[]
            )
            count = len(existing.get("ids", []))
            if count == 0:
                return ""
            results = self._col.query(
                query_texts=[text],
                n_results=min(n, count),
                where={"source": {"$in": paths}},
            )
            docs = results.get("documents", [[]])[0]
            return "\n\n---\n\n".join(d for d in docs if d)
        except Exception:
            logger.exception("Erreur query_filtered")
            return ""

    def query(self, text: str, n_results: Optional[int] = None) -> str:
        n = n_results if n_results is not None else self._n_results
        return self._query_lru(text, n)

    def query_filtered(self, text: str, paths: list, n_results: Optional[int] = None) -> str:
        if not paths:
            return ""
        n = n_results if n_results is not None else self._n_results
        return self._query_filtered_lru(text, tuple(sorted(paths)), n)

    def get_indexed_files(self) -> list:
        result = self._col.get(include=["metadatas"])
        sources = {m["source"] for m in result["metadatas"] if m and "source" in m}
        return sorted(sources)

    def describe_indexed_files(self) -> list[dict]:
        """Un descriptif par fichier indexé : chemins, chunks, mtime, indexation.

        Sert à choisir quoi retirer du corpus. Le **nombre de chunks** est la
        mesure qui compte : c'est ce que le fichier occupe réellement dans
        l'index, et pas sa taille sur le disque — un PDF de 40 Mo tout en images
        peut ne peser qu'un chunk. C'est aussi ce qui rend visible le symptôme du
        §3.3 bis : un fichier accepté que le moteur ne sait pas lire s'indexe à
        **zéro chunk**, en silence. Il n'apparaît alors pas ici du tout, faute de
        métadonnée — et c'est précisément l'information utile.

        ``indexé_le`` n'existe que pour ce qui a été indexé depuis son ajout :
        rien ne permet de le reconstituer pour l'existant, et le déduire du
        ``mtime`` du fichier serait faux (celui-ci dit quand le fichier a été
        modifié, pas quand on l'a lu). Absent = l'appelant affiche
        « non disponible ».

        Une seule lecture de l'index pour tous les fichiers : un appel par
        fichier relirait tout le magasin à chaque fois.
        """
        try:
            result = self._col.get(include=["metadatas"])
        except Exception:
            logger.exception("Erreur lecture de l'index pour le descriptif")
            return []

        par_source: dict[str, dict] = {}
        for m in result.get("metadatas", []) or []:
            if not m or "source" not in m:
                continue
            fiche = par_source.setdefault(
                m["source"], {"chemin": m["source"], "chunks": 0,
                              "mtime": m.get("mtime", 0.0), "indexé_le": ""})
            fiche["chunks"] += 1
            if not fiche["indexé_le"] and m.get("indexé_le"):
                fiche["indexé_le"] = m["indexé_le"]
        return sorted(par_source.values(), key=lambda f: f["chemin"])

    def remove_source(self, path: str) -> int:
        """Retire un fichier du corpus. Rend le nombre de chunks supprimés.

        ⚠️ **Ne touche pas au fichier sur le disque.** C'est un retrait de
        l'INDEX : le fichier reste là où il est et se réindexera si on le
        réimporte. Un bouton dans un panneau ne doit pas effacer les fichiers de
        quelqu'un, et l'irréversible ne doit jamais être le comportement par
        défaut d'une action nommée « supprimer » dans une liste.

        Les caches de requête sont invalidés, comme après une indexation : sans
        ça, `query()` continuerait de servir des extraits d'un fichier retiré —
        une réponse fondée sur un document que l'utilisateur croit supprimé.

        Rend ``0`` si le fichier n'était pas indexé ; l'appelant en fait un 404.
        """
        cible = str(path)
        avant = sum(1 for f in self.describe_indexed_files() if f["chemin"] == cible)
        if avant == 0:
            return 0
        self._col.delete(where={"source": cible})
        self._query_lru.cache_clear()
        self._query_filtered_lru.cache_clear()
        logger.info("Fichier retiré du corpus : %s (%d chunks)", cible, avant)
        return avant

    @staticmethod
    def read_file_text(path: str) -> str:
        return RAGEngine._extract_text_from_path(path)

    @staticmethod
    def read_pdf_text(path: str) -> str:
        """Backward-compat alias for read_file_text."""
        return RAGEngine._extract_text_from_path(path)

    def _indexed_mtimes(self) -> dict:
        """source → mtime déjà indexé (pour sauter les fichiers inchangés)."""
        try:
            result = self._col.get(include=["metadatas"])
        except Exception:
            logger.exception("Erreur lecture index pour scan incrémental")
            return {}
        out: dict[str, float] = {}
        for m in result.get("metadatas", []) or []:
            if m and "source" in m and "mtime" in m:
                out[m["source"]] = m["mtime"]
        return out

    def _initial_scan(self, folder: str) -> None:
        """Scan initial d'un dossier surveillé : (ré)indexe les fichiers nouveaux
        ou modifiés depuis la dernière fois. Tourne en tâche de fond (cf. watch).
        """
        indexed = self._indexed_mtimes()
        scanned = reindexed = 0
        for ext in SUPPORTED_EXTENSIONS:
            for file_path in Path(folder).rglob(f"*{ext}"):
                scanned += 1
                sp = str(file_path)
                try:
                    mtime = os.path.getmtime(sp)
                except OSError:
                    continue
                # Inchangé depuis la dernière indexation → pas de ré-embedding.
                if abs(indexed.get(sp, -1.0) - mtime) < 1e-6:
                    continue
                try:
                    self.index_file(sp)
                    reindexed += 1
                except Exception:
                    logger.exception("Erreur lors de l'indexation de %s", file_path)
        if reindexed:
            logger.info(
                "RAG scan initial %s : %d/%d fichier(s) (ré)indexé(s)",
                folder, reindexed, scanned,
            )

    def watch(self, folder: str) -> None:
        folder = str(folder)
        if not os.path.isdir(folder):
            return

        # L'observer démarre tout de suite (léger) : modifs/créations captées sans
        # attendre le scan initial.
        handler = _FileHandler(self)
        observer = Observer()
        observer.schedule(handler, folder, recursive=True)
        observer.daemon = True
        observer.start()

        # Scan initial en tâche de fond : il pouvait ré-embedder des dizaines de
        # fichiers (plusieurs minutes). Synchrone ici, il bloquait l'import de
        # core.runtime → uvicorn ne répondait pas et l'app restait figée sur
        # « Chargement… » au démarrage.
        threading.Thread(
            target=self._initial_scan, args=(folder,), daemon=True,
            name=f"rag-scan-{Path(folder).name}",
        ).start()


class _FileHandler(FileSystemEventHandler):
    def __init__(self, engine: RAGEngine):
        self._engine = engine

    def on_created(self, event):
        if not event.is_directory:
            ext = Path(event.src_path).suffix.lower()
            if ext in SUPPORTED_EXTENSIONS:
                try:
                    self._engine.index_file(event.src_path)
                except Exception:
                    logger.exception("Erreur lors de l'indexation de %s", event.src_path)
